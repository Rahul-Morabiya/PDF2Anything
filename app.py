from flask import Flask, request, send_from_directory, jsonify, render_template
import os
import uuid
from google.cloud import texttospeech
import requests
import fitz  # PyMuPDF for PDF text extraction
from pptx import Presentation
from flask_cors import CORS

# Setup Flask App
app = Flask(__name__)
os.makedirs("uploads", exist_ok=True)
os.makedirs("static", exist_ok=True)
os.makedirs("conversations", exist_ok=True)
os.makedirs("ppt", exist_ok=True)
os.makedirs("analysis", exist_ok=True)

# Load API keys and configurations
GEMINI_API_KEY = r"AIzaSyD7mrdGOTc-bccuNNPMbD254"  # Set in terminal using export GEMINI_API_KEY="your_api_key"
GOOGLE_CLOUD_TTS_KEY_PATH = r"C:\Users\Asus\OneDrive\Desktop\Projects\Cactus\bold-vortex-e3-c547f537c293.json"  # Path to service account JSON file

# Initialize Google Cloud TTS client
def initialize_google_tts_client():
    try:
        client = texttospeech.TextToSpeechClient.from_service_account_file(GOOGLE_CLOUD_TTS_KEY_PATH)
        return client
    except Exception as e:
        print("Error initializing TTS client:", e)
        raise

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        full_text = "\n".join(page.get_text("text") for page in doc)
        return full_text.strip()
    except Exception as e:
        print("Error extracting text from PDF:", e)
        return ""

# Generate podcast script using Google Gemini API
def generate_podcast_script(research_text, audience, tone):
    api_url = "https://generativelanguage.googleapis.com/v1/models/gemini-pro:generateContent"
    headers = {"Content-Type": "application/json"}
    
    prompt = (
        f"Audience: {audience.capitalize()}\n"
        f"Tone: {tone.capitalize()}\n\n"
        f"Create a conversational script discussing the following research paper:\n\n{research_text}"
    )
    
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    print("Calling Gemini API with prompt:", prompt[:100], "...")
    response = requests.post(f"{api_url}?key={GEMINI_API_KEY}", headers=headers, json=payload)
    
    if response.status_code == 200:
        result = response.json()
        if 'candidates' in result and result['candidates']:
            conversation = result['candidates'][0]['content']['parts'][0]['text']
            return conversation
        else:
            print("Gemini API response did not include candidates:", result)
    else:
        print("Gemini API request failed with status:", response.status_code, response.text)
    return "Error: Failed to generate podcast script."

# Generate audio from the podcast script using Google Cloud TTS
def generate_audio_from_text(client, text, voice_params, output_filename):
    synthesis_input = texttospeech.SynthesisInput(text=text)
    audio_config = texttospeech.AudioConfig(audio_encoding=texttospeech.AudioEncoding.MP3)
    
    try:
        response = client.synthesize_speech(input=synthesis_input, voice=voice_params, audio_config=audio_config)
        with open(output_filename, "wb") as out:
            out.write(response.audio_content)
    except Exception as e:
        print("Error generating audio:", e)
        raise

# Generate PPT from text (using python-pptx)
def generate_ppt_from_text(text, theme, num_slides):
    prs = Presentation()
    # Use a Title and Content layout (layout index 1)
    slide_layout = prs.slide_layouts[1]
    
    # Split the research text into roughly equal parts for the requested number of slides
    total_chars = len(text)
    if num_slides <= 0:
        num_slides = 1
    chunk_size = total_chars // num_slides

    for i in range(num_slides):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Slide {i+1}"
        
        start = i * chunk_size
        end = (i + 1) * chunk_size if i < num_slides - 1 else total_chars
        content_text = text[start:end].strip()
        
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content_text
        else:
            # Fallback: add a textbox if no placeholder available
            left = top = width = height = 100
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = content_text

    ppt_filename = f"ppt/{uuid.uuid4().hex}.pptx"
    prs.save(ppt_filename)
    return ppt_filename

# Handle file upload for Podcast
@app.route('/upload', methods=['POST'])
def upload_file():
    print("Received /upload request")
    print("Files received:", list(request.files.keys()))
    print("Form data received:", request.form)
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    # Retrieve customization options from form data
    audience = request.form.get('audience', 'scientific')
    length = request.form.get('length', '2min')
    tone = request.form.get('tone', 'formal')

    # Save uploaded PDF
    pdf_filename = f"uploads/{uuid.uuid4().hex}.pdf"
    try:
        file.save(pdf_filename)
    except Exception as e:
        print("Error saving uploaded file:", e)
        return jsonify({'error': 'Failed to save uploaded file'}), 500

    # Extract text from the PDF
    research_text = extract_text_from_pdf(pdf_filename)
    if not research_text:
        return jsonify({'error': 'Failed to extract text from PDF'}), 400

    # Generate the podcast script
    podcast_script = generate_podcast_script(research_text, audience, tone)
    if podcast_script.startswith("Error:"):
        return jsonify({'error': podcast_script}), 500

    # Initialize TTS client
    try:
        client = initialize_google_tts_client()
    except Exception as e:
        return jsonify({'error': 'TTS client initialization failed'}), 500

    # Generate audio from the podcast script
    audio_filename = f"static/{uuid.uuid4().hex}.mp3"
    voice_params = texttospeech.VoiceSelectionParams(
        language_code="en-US", ssml_gender=texttospeech.SsmlVoiceGender.FEMALE
    )
    try:
        generate_audio_from_text(client, podcast_script, voice_params, audio_filename)
    except Exception as e:
        return jsonify({'error': 'Failed to generate audio'}), 500

    # Save the conversation script as a text file
    conversation_filename = f"conversations/{uuid.uuid4().hex}.txt"
    try:
        with open(conversation_filename, 'w', encoding='utf-8') as f:
            f.write(podcast_script)
    except Exception as e:
        print("Error saving conversation script:", e)
        return jsonify({'error': 'Failed to save conversation script'}), 500

    # Return response with file paths for download
    return jsonify({
        "audio_url": os.path.basename(audio_filename),
        "conversation_url": os.path.basename(conversation_filename)
    })

# Handle file upload for PPT Generation
@app.route('/upload_ppt', methods=['POST'])
def upload_ppt_file():
    print("Received /upload_ppt request")
    print("Files received:", list(request.files.keys()))
    print("Form data received:", request.form)
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    num_slides = int(request.form.get('num_slides', 5))
    audience = request.form.get('audience', 'scientific')
    theme = request.form.get('theme', 'default')

    pdf_filename = f"uploads/{uuid.uuid4().hex}.pdf"
    try:
        file.save(pdf_filename)
    except Exception as e:
        print("Error saving uploaded PDF for PPT:", e)
        return jsonify({'error': 'Failed to save uploaded PDF'}), 500

    research_text = extract_text_from_pdf(pdf_filename)
    if not research_text:
        return jsonify({'error': 'Failed to extract text from PDF'}), 400

    try:
        ppt_filename = generate_ppt_from_text(research_text, theme, num_slides)
    except Exception as e:
        print("Error generating PPT:", e)
        return jsonify({'error': 'Failed to generate PPT'}), 500

    return jsonify({"ppt_url": os.path.basename(ppt_filename)})

# Function for analysis using Gemini API
def gemini_api_analysis(research_text, analysis_type):
    api_url = "https://generativelanguage.googleapis.com/v1/models/gemini-pro:generateContent"
    headers = {"Content-Type": "application/json"}
    
    prompt = (
        f"Analysis Type: {analysis_type.capitalize()}\n\n"
        f"Analyze the following research paper:\n\n{research_text}"
    )
    
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    print("Calling Gemini API for analysis with prompt:", prompt[:100], "...")
    response = requests.post(f"{api_url}?key={GEMINI_API_KEY}", headers=headers, json=payload)
    
    if response.status_code == 200:
        result = response.json()
        if 'candidates' in result and result['candidates']:
            analysis = result['candidates'][0]['content']['parts'][0]['text']
            return analysis
        else:
            print("Gemini API analysis response did not include candidates:", result)
    else:
        print("Gemini API analysis request failed with status:", response.status_code, response.text)
    return "Error: Failed to generate analysis."

# Endpoint to handle file upload and analysis request
@app.route('/upload_analysis', methods=['POST'])
def upload_analysis():
    print("Received /upload_analysis request")
    print("Files received:", list(request.files.keys()))
    print("Form data received:", request.form)
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    analysis_type = request.form.get('analysis_type', 'summary')
    pdf_filename = f"uploads/{uuid.uuid4().hex}.pdf"
    try:
        file.save(pdf_filename)
    except Exception as e:
        print("Error saving uploaded PDF for analysis:", e)
        return jsonify({'error': 'Failed to save uploaded PDF'}), 500

    research_text = extract_text_from_pdf(pdf_filename)
    if not research_text:
        return jsonify({'error': 'Failed to extract text from PDF'}), 400

    analysis_result = gemini_api_analysis(research_text, analysis_type)
    if analysis_result.startswith("Error:"):
        return jsonify({'error': analysis_result}), 500

    analysis_filename = f"analysis/{uuid.uuid4().hex}.txt"
    try:
        with open(analysis_filename, 'w', encoding='utf-8') as f:
            f.write(analysis_result)
    except Exception as e:
        print("Error saving analysis result:", e)
        return jsonify({'error': 'Failed to save analysis result'}), 500

    return jsonify({
        'analysis_url': os.path.basename(analysis_filename),
        'analysis_result': analysis_result
    })

# Serve static files for download
@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(os.path.join(app.root_path, 'static'), filename)

@app.route('/download/conversation/<filename>')
def download_conversation(filename):
    return send_from_directory(os.path.join(app.root_path, 'conversations'), filename)

@app.route('/download/ppt/<filename>')
def download_ppt(filename):
    return send_from_directory(os.path.join(app.root_path, 'ppt'), filename)

@app.route('/download/analysis/<filename>')
def download_analysis(filename):
    return send_from_directory(os.path.join(app.root_path, 'analysis'), filename)

# Serve the index.html from templates
@app.route('/')
def home():
    return render_template('index.html')

if __name__ == '__main__':
    CORS(app)  # Allow cross-origin requests
    app.run(debug=True)
